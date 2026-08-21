import os
import io
import fitz  # PyMuPDF
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

def hex_to_rgb(hex_color):
    if isinstance(hex_color, int):
        r = (hex_color >> 16) & 0xFF
        g = (hex_color >> 8) & 0xFF
        b = hex_color & 0xFF
        return RGBColor(r, g, b)
    return RGBColor(0, 0, 0)

def parse_page_spec(spec_str, total_pages):
    """
    Parse a string like '2', '3, 5-7' into a set of 0-indexed page numbers.
    """
    if not spec_str or not spec_str.strip():
        return set()
    
    pages = set()
    parts = spec_str.split(',')
    for part in parts:
        part = part.strip()
        if not part:
            continue
        if '-' in part:
            subparts = part.split('-')
            if len(subparts) == 2 and subparts[0].strip().isdigit() and subparts[1].strip().isdigit():
                start = int(subparts[0].strip())
                end = int(subparts[1].strip())
                for p in range(min(start, end), max(start, end) + 1):
                    if 1 <= p <= total_pages:
                        pages.add(p - 1)
        elif part.isdigit():
            p = int(part)
            if 1 <= p <= total_pages:
                pages.add(p - 1)
    return pages

class PDFToPPTConverter:
    def __init__(self, progress_callback=None):
        self.progress_callback = progress_callback

    def log(self, message):
        if self.progress_callback:
            self.progress_callback(message)
        else:
            print(message)

    def convert(self, pdf_path, ppt_path, select_pages_str="", ignore_pages_str="", custom_height_px=None):
        self.log(f"Opening PDF: {os.path.basename(pdf_path)}...")
        pdf_doc = fitz.open(pdf_path)
        total_pages = len(pdf_doc)

        if total_pages == 0:
            raise ValueError("PDF file contains no pages.")

        # Open pdfplumber document for visual layout inspection
        plumber_doc = None
        try:
            plumber_doc = pdfplumber.open(pdf_path)
        except Exception as e:
            self.log(f"Notice: pdfplumber fallback disabled: {e}")

        # Determine target pages
        if select_pages_str and select_pages_str.strip():
            selected_set = parse_page_spec(select_pages_str, total_pages)
            if not selected_set:
                if total_pages == 1:
                    self.log(f"Notice: Selection '{select_pages_str}' out of range for 1-page PDF. Defaulting to Page 1.")
                    selected_set = {0}
                else:
                    raise ValueError(f"Selected pages '{select_pages_str}' are out of range (PDF has {total_pages} page(s)). Please specify pages between 1 and {total_pages}.")
        else:
            selected_set = set(range(total_pages))

        if ignore_pages_str and ignore_pages_str.strip():
            ignored_set = parse_page_spec(ignore_pages_str, total_pages)
        else:
            ignored_set = set()

        final_pages = [p for p in range(total_pages) if p in selected_set and p not in ignored_set]

        if not final_pages:
            raise ValueError(f"No pages left to convert after applying filters. Selected: {selected_set}, Ignored: {ignored_set} (PDF Total: {total_pages} pages).")

        # Prepare separate presentations for normal (horizontal/standard) and vertical slides
        prs_normal = Presentation()
        prs_vertical = Presentation()

        normal_count = 0
        vertical_count = 0

        # Set default presentation dimensions based on first selected page or custom height
        first_rect = pdf_doc[final_pages[0]].rect
        def_w_in = Inches(first_rect.width / 72.0)
        
        if custom_height_px and custom_height_px > 0:
            def_h_in = Inches(custom_height_px / 96.0) # 96 px = 1 inch
        else:
            def_h_in = Inches(first_rect.height / 72.0)

        prs_normal.slide_width = def_w_in
        prs_normal.slide_height = def_h_in

        prs_vertical.slide_width = def_w_in
        prs_vertical.slide_height = def_h_in

        # Configure Slide Master theme colors to prevent PPT from defaulting hyperlinks to Blue
        for p_master in [prs_normal, prs_vertical]:
            try:
                slide_master = p_master.slide_master
                theme_elem = slide_master.element.xpath('//a:clrScheme')
                if theme_elem:
                    clr_scheme = theme_elem[0]
                    from pptx.oxml.xmlchemy import OxmlElement
                    # Find or modify hlink
                    for child in list(clr_scheme):
                        if child.tag.endswith('hlink'):
                            for c in list(child):
                                child.remove(c)
                            srgbClr = OxmlElement('a:srgbClr')
                            srgbClr.set('val', 'E50000') # Red theme link
                            child.append(srgbClr)
            except Exception as e:
                pass

        blank_layout_normal = prs_normal.slide_layouts[6]
        blank_layout_vertical = prs_vertical.slide_layouts[6]

        for page_num in final_pages:
            self.log(f"Processing Page {page_num + 1}/{total_pages}...")
            page = pdf_doc[page_num]
            rect = page.rect

            # Extract visual character map from pdfplumber for current page if available
            plumber_chars = []
            if plumber_doc and page_num < len(plumber_doc.pages):
                try:
                    plumber_chars = plumber_doc.pages[page_num].chars
                except Exception:
                    pass

            # Check if page is vertical (portrait: height > width)
            is_vertical = rect.height > rect.width

            if is_vertical:
                prs = prs_vertical
                blank_layout = blank_layout_vertical
                vertical_count += 1
            else:
                prs = prs_normal
                blank_layout = blank_layout_normal
                normal_count += 1

            slide = prs.slides.add_slide(blank_layout)

            # Adjust slide height if custom height specified or page specific
            if custom_height_px and custom_height_px > 0:
                prs.slide_height = Inches(custom_height_px / 96.0)
            
            # 1. Detect and crop isolated vector drawing clusters (ratings, banners, boxes, icons)
            try:
                drawings = page.get_drawings()
                drawing_rects = []
                for draw in drawings:
                    r = draw.get("rect")
                    if r and r.width > 0 and r.height > 0:
                        if r.width >= rect.width * 0.9 and r.height >= rect.height * 0.9:
                            continue
                        drawing_rects.append(r)

                clean_doc = fitz.open(pdf_path)
                clean_page = clean_doc[page_num]
                for block in clean_page.get_text("blocks"):
                    if block[4].strip():
                        clean_page.add_redact_annot(fitz.Rect(block[:4]), fill=None)
                clean_page.apply_redactions()

                for d_rect in drawing_rects:
                    pix = clean_page.get_pixmap(clip=d_rect, dpi=150)
                    img_bytes = pix.tobytes("png")
                    img_stream = io.BytesIO(img_bytes)

                    left = Inches(d_rect.x0 / 72.0)
                    top = Inches(d_rect.y0 / 72.0)
                    width = Inches(d_rect.width / 72.0)
                    height = Inches(d_rect.height / 72.0)

                    try:
                        slide.shapes.add_picture(img_stream, left, top, width, height)
                    except Exception as e:
                        pass
                clean_doc.close()
            except Exception as e:
                self.log(f"Notice extracting vector graphics on page {page_num+1}: {e}")

            # 2. Extract embedded raster bitmap images
            try:
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    image_rects = page.get_image_rects(xref)
                    for r in image_rects:
                        if r.width >= rect.width * 0.9 and r.height >= rect.height * 0.9:
                            continue
                        
                        pix = page.get_pixmap(clip=r, dpi=200, colorspace=fitz.csRGB)
                        image_bytes = pix.tobytes("png")
                        
                        ras_stream = io.BytesIO(image_bytes)
                        left = Inches(r.x0 / 72.0)
                        top = Inches(r.y0 / 72.0)
                        width = Inches(r.width / 72.0)
                        height = Inches(r.height / 72.0)
                        slide.shapes.add_picture(ras_stream, left, top, width, height)
            except Exception as e:
                self.log(f"Warning extracting bitmap image on page {page_num+1}: {e}")

            # Extract hyperlink targets on current page
            links = page.get_links()

            # 3. Extract and overlay 100% editable text blocks with visual pdfplumber analysis
            text_page = page.get_text("dict")
            for block in text_page.get("blocks", []):
                if block.get("type") == 0:  # Text block
                    bbox = block.get("bbox")
                    b_x0, b_y0, b_x1, b_y1 = bbox
                    
                    box_width = (b_x1 - b_x0) + 40.0
                    box_height = (b_y1 - b_y0) + 10.0
                    
                    txBox = slide.shapes.add_textbox(
                        Inches(b_x0 / 72.0),
                        Inches(b_y0 / 72.0),
                        Inches(box_width / 72.0),
                        Inches(box_height / 72.0)
                    )
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0)
                    tf.margin_right = Inches(0)
                    tf.margin_top = Inches(0)
                    tf.margin_bottom = Inches(0)

                    all_spans = []
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            if span.get("text", ""):
                                all_spans.append(span)

                    if not all_spans:
                        continue

                    # -------------------------------------------------------------------------
                    # Hybrid Visual Line Grouping
                    # Group spans by visual Y-center tolerance (3.5pt) so list numbers & superscripts
                    # remain on the exact same paragraph line as main text.
                    # -------------------------------------------------------------------------
                    lines_grouped = []
                    for span in all_spans:
                        s_bbox = span.get("bbox")
                        s_y_mid = (s_bbox[1] + s_bbox[3]) / 2.0
                        
                        matched_group = None
                        for group in lines_grouped:
                            g_y_mid = group["y_mid"]
                            if abs(s_y_mid - g_y_mid) <= 3.5:
                                matched_group = group
                                break
                        
                        if matched_group:
                            matched_group["spans"].append(span)
                        else:
                            lines_grouped.append({"y_mid": s_y_mid, "spans": [span]})

                    lines_grouped.sort(key=lambda g: g["y_mid"])

                    for group in lines_grouped:
                        max_sz = max(sp.get("size", 12) for sp in group["spans"])
                        dominant_spans = [sp for sp in group["spans"] if sp.get("size", 12) >= max_sz * 0.88]
                        dom_baselines = [sp.get("origin", (0, sp.get("bbox")[3]))[1] for sp in dominant_spans]
                        group["dominant_baseline_y"] = sum(dom_baselines) / len(dom_baselines) if dom_baselines else group["y_mid"]
                        group["max_size"] = max_sz

                    # Extract all horizontal vector underline paths on page
                    underline_rects = []
                    try:
                        for draw in page.get_drawings():
                            for item in draw.get("items", []):
                                if item[0] in ("l", "re"):
                                    p0 = item[1]
                                    p1 = item[2] if len(item) > 2 else p0
                                    y0 = min(p0.y, p1.y) if hasattr(p0, 'y') else draw["rect"].y0
                                    y1 = max(p0.y, p1.y) if hasattr(p0, 'y') else draw["rect"].y1
                                    if abs(y1 - y0) <= 3.0:
                                        underline_rects.append(draw["rect"])
                    except Exception:
                        pass

                    first_paragraph = True
                    for group in lines_grouped:
                        g_spans = sorted(group["spans"], key=lambda s: s.get("bbox")[0])
                        
                        if first_paragraph:
                            p = tf.paragraphs[0]
                            first_paragraph = False
                        else:
                            p = tf.add_paragraph()

                        for idx, span in enumerate(g_spans):
                            text = span.get("text", "")
                            if not text:
                                continue
                            
                            if idx > 0 and not text.startswith(" ") and not text.startswith(",") and not text.startswith("."):
                                text = " " + text

                            run = p.add_run()
                            run.text = text
                            
                            # Font styling
                            size = span.get("size", 12)
                            font_name = span.get("font", "Calibri")
                            if "+" in font_name:
                                font_name = font_name.split("+")[-1]
                            run.font.name = font_name

                            # Flags: 1=superscript, 2=italic, 4=serifed, 8=monospaced, 16=bold
                            flags = span.get("flags", 0)
                            if flags & 2 or "italic" in font_name.lower():
                                run.font.italic = True
                            if flags & 16 or "bold" in font_name.lower():
                                run.font.bold = True

                            # Underline detection
                            s_bbox = span.get("bbox")
                            s_rect = fitz.Rect(s_bbox)
                            has_underline = bool("underline" in font_name.lower())
                            
                            if not has_underline:
                                for u_rect in underline_rects:
                                    if abs(u_rect.y0 - s_rect.y1) <= 3.5:
                                        overlap_x0 = max(u_rect.x0, s_rect.x0)
                                        overlap_x1 = min(u_rect.x1, s_rect.x1)
                                        overlap_w = max(0, overlap_x1 - overlap_x0)
                                        span_w = s_rect.width
                                        if span_w > 0 and (overlap_w / span_w) >= 0.5:
                                            if text.endswith(".") and u_rect.x1 < s_rect.x1 - 2.0:
                                                pass
                                            else:
                                                has_underline = True
                                                break
                            
                            if has_underline:
                                run.font.underline = True

                            # -----------------------------------------------------------------
                            # VISUAL SUPERSCRIPT & SUBSCRIPT ENGINE (Powered by pdfplumber + PyMuPDF)
                            # -----------------------------------------------------------------
                            span_origin_y = span.get("origin", (0, s_bbox[3]))[1]
                            dom_baseline_y = group["dominant_baseline_y"]
                            max_line_size = group["max_size"]

                            is_super = bool(flags & 1)
                            is_sub = False

                            baseline_shift = dom_baseline_y - span_origin_y
                            is_list_number = (idx == 0 and text.strip().replace('.', '').isdigit())

                            # Cross-verify with pdfplumber visual character layout engine if available
                            plumber_visual_super = False
                            plumber_visual_sub = False
                            if plumber_chars and not is_list_number:
                                # Find matching characters in pdfplumber for this span
                                matching_c = []
                                for c in plumber_chars:
                                    if abs(c["x0"] - s_bbox[0]) <= 2.0 and abs(c["top"] - s_bbox[1]) <= 2.5:
                                        matching_c.append(c)
                                if matching_c:
                                    c_top = sum(c["top"] for c in matching_c) / len(matching_c)
                                    # Compare against line top position
                                    line_top = s_bbox[1]
                                    if c["size"] < max_line_size * 0.88:
                                        if baseline_shift >= 0.8:
                                            plumber_visual_super = True
                                        elif baseline_shift <= -0.8:
                                            plumber_visual_sub = True

                            if not is_list_number and size <= max_line_size * 0.88:
                                if baseline_shift >= 1.0 or plumber_visual_super or (is_super and baseline_shift > -0.5):
                                    is_super = True
                                    is_sub = False
                                elif baseline_shift <= -1.0 or plumber_visual_sub:
                                    is_sub = True
                                    is_super = False

                            if is_super and not is_list_number:
                                run.font.superscript = True
                                run.font.size = Pt(max_line_size)
                            elif is_sub and not is_list_number:
                                run.font.subscript = True
                                run.font.size = Pt(max_line_size)
                            else:
                                run.font.size = Pt(size)

                            color = span.get("color", 0)
                            rgb_obj = hex_to_rgb(color)
                            run.font.color.rgb = rgb_obj

                            # Check if span falls within a PDF hyperlink region
                            s_rect = fitz.Rect(s_bbox)
                            for link in links:
                                l_rect = link.get("from")
                                if l_rect and s_rect.intersects(l_rect):
                                    uri = link.get("uri")
                                    if uri:
                                        try:
                                            txBox.click_action.hyperlink.address = uri
                                        except Exception:
                                            pass
                                    break

        if plumber_doc:
            try:
                plumber_doc.close()
            except Exception:
                pass

        saved_files = []
        base_no_ext, ext = os.path.splitext(ppt_path)

        if normal_count > 0 and vertical_count == 0:
            prs_normal.save(ppt_path)
            saved_files.append(ppt_path)
            self.log(f"Saved presentation: {os.path.basename(ppt_path)}")

        elif vertical_count > 0 and normal_count == 0:
            prs_vertical.save(ppt_path)
            saved_files.append(ppt_path)
            self.log(f"Saved presentation: {os.path.basename(ppt_path)}")

        else:
            prs_normal.save(ppt_path)
            saved_files.append(ppt_path)
            self.log(f"Saved standard presentation: {os.path.basename(ppt_path)}")

            vert_ppt_path = f"{base_no_ext}_vertical{ext}"
            prs_vertical.save(vert_ppt_path)
            saved_files.append(vert_ppt_path)
            self.log(f"Saved vertical presentation: {os.path.basename(vert_ppt_path)}")

        return saved_files

