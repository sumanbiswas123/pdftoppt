import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Hello Editable World!"
run.font.size = Pt(24)
run.font.bold = True
run.font.color.rgb = RGBColor(0, 102, 204)

output_path = "sample_test.pptx"
prs.save(output_path)
print("SUCCESS: Sample test PPTX created successfully.")
